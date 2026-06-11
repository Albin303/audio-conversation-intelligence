"""
Extract all visible text from a PowerPoint file into a plain-text report.

Usage (from project root):
    python scripts/extract_ppt_text.py

Dependencies: python-pptx (see requirements.txt)
"""

from __future__ import annotations

import argparse
from pathlib import Path

from pptx import Presentation


def iter_slide_lines(prs: Presentation):
    """Yield (slide_index, lines) for each slide."""
    for idx, slide in enumerate(prs.slides, start=1):
        lines = []
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            for paragraph in shape.text_frame.paragraphs:
                text = paragraph.text.strip()
                if text:
                    lines.append(text)
        yield idx, lines


def write_report(prs: Presentation, output: Path, source: Path) -> None:
    with output.open("w", encoding="utf-8") as f:
        f.write("# AI-Driven Customer Conversation Intelligence Platform\n\n")
        f.write("## Source\n")
        f.write(f"- Converted from `{source.name}` ({len(prs.slides)} slides)\n\n")
        f.write("## Slide Contents\n\n")

        for slide_idx, lines in iter_slide_lines(prs):
            f.write(f"--- Slide {slide_idx} ---\n")
            if lines:
                for line in lines:
                    f.write(f"- {line}\n")
            else:
                f.write("- (no text found)\n")
            f.write("\n")


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Extract text from a PowerPoint file")
    parser.add_argument(
        "--input",
        type=Path,
        default=Path("docs/assets/sales_conversation_intelligence_presentation.pptx"),
        help="PPTX file",
    )
    parser.add_argument(
        "--output",
        type=Path,
        default=Path("docs/assets/sales_conversation_intelligence_presentation.txt"),
        help="Destination text file (overwrite if exists)",
    )
    return parser.parse_args()


def main() -> None:
    args = parse_args()
    prs = Presentation(args.input)
    args.output.parent.mkdir(parents=True, exist_ok=True)
    write_report(prs, args.output, args.input)
    print(f"Extracted text from {args.input} -> {args.output}")


if __name__ == "__main__":
    main()
